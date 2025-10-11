#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
app_2.py — fixed, de-duplicated, resilient.

- Single robust core import (finance_core / Monthly_pipeline_3 / core.finance_core)
- Safe aliases for core symbols + shims
- No duplicate blocks, no shadowed names
- Guaranteed run_analysis (uses core if present, else wrapper)
"""

from __future__ import annotations
import io, os, sys, calendar, datetime as dt, tempfile
from dataclasses import dataclass
from typing import Dict, List, Optional, Tuple
from types import SimpleNamespace

import numpy as np
import pandas as pd
import streamlit as st

# ========================= Robust core import =========================
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
if BASE_DIR not in sys.path:
    sys.path.insert(0, BASE_DIR)

def _import_core():
    for name in ("finance_core", "Monthly_pipeline_3", "core.finance_core"):
        try:
            return __import__(name, fromlist=["*"])
        except Exception:
            continue
    raise ImportError("Could not import finance_core (also tried Monthly_pipeline_3, core.finance_core).")

try:
    _core = _import_core()
except Exception as e:
    st.error(
        "Could not import the finance core module.\n\n"
        "Ensure **finance_core.py** or **Monthly_pipeline_3.py** is in the SAME folder as this app."
    )
    st.code(
        "cwd: " + os.getcwd() + "\n"
        "here: " + BASE_DIR + "\n"
        "files: " + ", ".join(sorted(os.listdir(BASE_DIR))) + "\n"
        "sys.path[0..5]: " + ", ".join(map(str, sys.path[:6]))
    )
    st.exception(e)
    st.stop()

# ---- Optional/required core symbols with safe fallbacks ----
read_excel_any         = getattr(_core, "read_excel_any", None)
export_excel_core      = getattr(_core, "export_excel", None)
process_dataframe_core = getattr(_core, "process_dataframe", None)

translate_columns      = getattr(_core, "translate_columns", None)
normalize_core         = getattr(_core, "normalize", None)           # NOTE: alias name
compute_g1_transport   = getattr(_core, "compute_g1_transport", None)

validate_expenditures  = getattr(_core, "validate_expenditures", None)
normalize_expenditures = getattr(_core, "normalize_expenditures", None)
compare_expenditures   = getattr(_core, "compare_expenditures", None)
summarize_expenditures = getattr(_core, "summarize_expenditures", None)

build_pptx_core        = getattr(_core, "build_pptx", None)
build_pdf_core         = getattr(_core, "build_pdf", None)
load_month_amount_file_core = getattr(_core, "load_month_amount_file", None)

run_analysis_core      = getattr(_core, "run_analysis", None)

# Constants
CORRESPONDENT_MAP_DEFAULT = getattr(_core, "CORRESPONDENT_MAP_DEFAULT", {})
SPECIAL_CORR_DEFAULT      = getattr(_core, "SPECIAL_CORR_DEFAULT", [])
VAT_RATE_DEFAULT          = float(getattr(_core, "VAT_RATE", 0.12))
VAT_MODE_DEFAULT          = str(getattr(_core, "VAT_MODE", "extract"))
DATE_SOURCE_DEFAULT       = str(getattr(_core, "DATE_SOURCE", "Data of Document"))
TOP_N_NAMES_DEFAULT       = int(getattr(_core, "TOP_N_NAMES", 50))

# ---- Minimal fallbacks ----
if read_excel_any is None:
    def read_excel_any(path_or_buffer):
        try:
            return pd.read_excel(path_or_buffer)
        except Exception:
            return pd.read_excel(path_or_buffer, engine="xlrd")

# Export wrapper -> bytes for Streamlit
def export_excel_bytes(sheets_dict: Dict[str, pd.DataFrame]) -> bytes:
    if export_excel_core is None:
        # Local simple writer if core exporter is missing
        bio = io.BytesIO()
        with pd.ExcelWriter(bio, engine="xlsxwriter") as writer:
            for name, df in sheets_dict.items():
                df.to_excel(writer, sheet_name=name[:31], index=False)
        bio.seek(0)
        return bio.read()

    out_path = os.path.join(
        tempfile.gettempdir(),
        f"financial_report__{dt.datetime.now().strftime('%Y%m%d-%H%M%S')}.xlsx"
    )
    export_excel_core(sheets_dict, out_path)
    with open(out_path, "rb") as f:
        return f.read()

# ============================ Local shims =============================

@dataclass
class _ValResult:
    ok: bool
    missing: List[str]
    suggestions: str = ""

REQUIRED_COLS = [
    "Correspondent", "Number of Documents", "Data of Document", "Data of transaction",
    "Number of request", "Material/SAP Code", "Name", "Qty", "Measurement",
    "Amount", "Currency", "Warranty"
]

# validate_revenue (accepts RU→EN via translate_columns if available)
def validate_revenue(df: pd.DataFrame) -> _ValResult:
    try:
        d = translate_columns(df, verbose=False) if translate_columns else df
    except Exception as e:
        return _ValResult(False, ["(translate_columns failed)"], str(e))
    missing = [c for c in REQUIRED_COLS if c not in d.columns]
    if missing:
        return _ValResult(False, missing,
                          "Make sure the SAP file headers (RU/EN) translate into: " + ", ".join(REQUIRED_COLS))
    return _ValResult(True, [])

# normalize_revenue (translate → normalize_core → compute_g1_transport)
def normalize_revenue(df: pd.DataFrame,
                      corr_map: Dict[int, str] | None = None,
                      special_corr: List[int] | None = None) -> pd.DataFrame:
    d = translate_columns(df, verbose=False) if translate_columns else df
    if normalize_core:
        d, _ = normalize_core(d, corr_map or CORRESPONDENT_MAP_DEFAULT)
    if compute_g1_transport:
        d = compute_g1_transport(d, special_corr or SPECIAL_CORR_DEFAULT)
    return d

def months_between(start_ym: str, end_ym: str) -> List[str]:
    sy, sm = map(int, start_ym.split("-")); ey, em = map(int, end_ym.split("-"))
    cur = dt.date(sy, sm, 1); end = dt.date(ey, em, 1); out = []
    while cur <= end:
        out.append(f"{cur.year:04d}-{cur.month:02d}")
        cur = dt.date(cur.year + (cur.month == 12), (cur.month % 12) + 1, 1)
    return out

def _calc_after_vat(amount_vat_incl: float, vat_rate: float, vat_mode: str) -> float:
    vm = (vat_mode or "extract").lower()
    return amount_vat_incl / (1.0 + float(vat_rate)) if vm == "extract" else amount_vat_incl * (1.0 - float(vat_rate))

def _pick_date_column(df: pd.DataFrame, date_cols: List[str]) -> str:
    for c in date_cols:
        if c in df.columns: return c
    for c in df.columns:
        if pd.api.types.is_datetime64_any_dtype(df[c]): return c
    raise ValueError(f"No usable date columns found. Tried: {date_cols}")

def _ensure_month_key(df: pd.DataFrame, date_col: str) -> pd.DataFrame:
    d = df.copy()
    d[date_col] = pd.to_datetime(d[date_col], errors="coerce")
    d["Month"] = d[date_col].dt.to_period("M").astype(str)
    return d

def _subset_by_months(df: pd.DataFrame, months: List[str]) -> pd.DataFrame:
    return df[df["Month"].isin(months)].copy()

def _days_in_month(ym: str) -> int:
    y, m = map(int, ym.split("-"))
    return calendar.monthrange(y, m)[1]

def _active_days_factor(df_month: pd.DataFrame,
                        date_col: str,
                        vat_rate: float,
                        vat_mode: str,
                        nonempty_only: bool,
                        exclude_sundays: bool) -> Tuple[int, int, float]:
    if df_month.empty:
        return 0, 0, 0.0
    month_key = str(df_month["Month"].iloc[0]); y, m = map(int, month_key.split("-"))
    month_days = _days_in_month(month_key)

    d = df_month.copy()
    d[date_col] = pd.to_datetime(d[date_col], errors="coerce")
    d["_date"] = d[date_col].dt.date
    g = d.groupby("_date", dropna=False).agg(gross=("Amount", "sum"), g1=("g1_transport", "sum")).reset_index()
    g["after_vat_excl_cc"] = _calc_after_vat(g["gross"] - g["g1"], vat_rate, vat_mode)

    def is_sun(d0: dt.date) -> bool:
        try: return dt.date(d0.year, d0.month, d0.day).weekday() == 6
        except Exception: return False

    mask = (g["after_vat_excl_cc"] > 0) if nonempty_only else pd.Series(True, index=g.index)
    if exclude_sundays:
        mask = mask & (~g["_date"].apply(is_sun))

    active_days = int(mask.sum())

    if exclude_sundays:
        month_days = sum(1 for d0 in pd.date_range(dt.date(y, m, 1), dt.date(y, m, month_days)) if d0.weekday() != 6)

    factor = (month_days / active_days) if active_days else 0.0
    return active_days, month_days, factor

def _read_any_cached(file_bytes: bytes, filename: str) -> pd.DataFrame:
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".csv":
        return pd.read_csv(io.BytesIO(file_bytes))
    return read_excel_any(io.BytesIO(file_bytes))

# ---- Month/Amount loader (local tolerant) ----
def parse_month_ru(s: str):
    # Minimal helper: returns (year, month_num, name, "YYYY-MM") if parseable else None
    months = {
        "янв":1,"фев":2,"мар":3,"апр":4,"май":5,"мая":5,"июн":6,"июл":7,"авг":8,"сен":9,"сент":9,"окт":10,"ноя":11,"дек":12
    }
    ss = s.strip().lower()
    for key, m in months.items():
        if key in ss:
            nums = [int(t) for t in "".join(ch if ch.isdigit() or ch==" " else " " for ch in ss).split() if t.isdigit()]
            y = next((n for n in nums if n>=2000), None)
            if y:
                return (y, m, key, f"{y:04d}-{m:02d}")
    # ISO-like?
    try:
        y,m = ss.split("-")[:2]
        return (int(y), int(m), "", f"{int(y):04d}-{int(m):02d}")
    except Exception:
        return (None, None, None, None)

def load_month_amount_file_local(file_like_or_bytes) -> pd.DataFrame:
    if isinstance(file_like_or_bytes, (bytes, bytearray)):
        bio = io.BytesIO(file_like_or_bytes)
        try:
            df = pd.read_excel(bio)
        except Exception:
            bio.seek(0); df = pd.read_csv(bio)
    else:
        try:
            df = pd.read_excel(file_like_or_bytes)
        except Exception:
            df = pd.read_csv(file_like_or_bytes)

    cols = {str(c).strip(): c for c in df.columns}
    month_col = next((cols[c] for c in ("Month","Месяц") if c in cols), None)
    amt_col   = next((cols[c] for c in ("Amount_USD","USD","Amount") if c in cols), None)

    if month_col is None:
        for c in df.columns:
            if "month" in str(c).lower(): month_col = c; break
    if amt_col is None:
        for c in df.columns:
            if pd.api.types.is_numeric_dtype(df[c]): amt_col = c; break

    if month_col is None or amt_col is None:
        return pd.DataFrame({"Month": [], "Amount_USD": []})

    out_month = []
    for v in df[month_col].astype(str).tolist():
        _, _, _, ym = parse_month_ru(v)
        out_month.append(ym if ym else v.strip())

    amounts = pd.to_numeric(df[amt_col], errors="coerce").fillna(0.0)

    res = pd.DataFrame({"Month": out_month, "Amount_USD": amounts})
    def _fix_ym(s):
        try:
            y, m = s.split("-"); return f"{int(y):04d}-{int(m):02d}"
        except Exception:
            return s
    res["Month"] = res["Month"].astype(str).str.strip().apply(_fix_ym)
    return res.groupby("Month", as_index=False)["Amount_USD"].sum()

def _cc_amount_for_month(cc_df: Optional[pd.DataFrame], ym: str) -> float:
    if cc_df is None or cc_df.empty: return 0.0
    if "Month" not in cc_df.columns or "Amount_USD" not in cc_df.columns: return 0.0
    s = cc_df.loc[cc_df["Month"].astype(str) == ym, "Amount_USD"]
    return float(s.iloc[0]) if not s.empty and pd.notna(s.iloc[0]) else 0.0

def compare_ranges_revenue(
    df_rev: pd.DataFrame,
    ar_start: str, ar_end: str,
    pr_start: str, pr_end: str,
    date_cols: List[str],
    vat_rate: float, vat_mode: str,
    cc_actual_df: Optional[pd.DataFrame],
    cc_prev_df: Optional[pd.DataFrame],
    forecast_last_ym: Optional[str] = None,
    nonempty_only: bool = True,
    exclude_sundays: bool = True
) -> Tuple[pd.DataFrame, pd.DataFrame]:
    if df_rev is None or df_rev.empty:
        return pd.DataFrame(), pd.DataFrame()

    date_col = _pick_date_column(df_rev, date_cols)
    d = _ensure_month_key(df_rev, date_col)

    actual_months   = months_between(ar_start, ar_end)
    previous_months = months_between(pr_start, pr_end)

    d_actual = _subset_by_months(d, actual_months)
    d_prev   = _subset_by_months(d, previous_months)

    def _per_month_after_vat(df_in: pd.DataFrame) -> pd.DataFrame:
        if df_in.empty:
            return pd.DataFrame({"Month": [], "AfterVAT_excl_CC": []})
        g = df_in.groupby("Month", as_index=False).agg(Gross=("Amount", "sum"),
                                                      G1=("g1_transport", "sum") if "g1_transport" in df_in.columns else ("Amount","sum"))
        if "g1_transport" not in g.columns:
            g["G1"] = 0.0
        g["AfterVAT_excl_CC"] = _calc_after_vat(g["Gross"] - g["G1"], vat_rate, vat_mode)
        return g[["Month", "AfterVAT_excl_CC"]]

    act_pm  = _per_month_after_vat(d_actual)
    prev_pm = _per_month_after_vat(d_prev)

    if forecast_last_ym and (forecast_last_ym in actual_months):
        last_m_df = d_actual[d_actual["Month"] == forecast_last_ym]
        if not last_m_df.empty:
            _, _, factor = _active_days_factor(last_m_df, date_col, vat_rate, vat_mode, nonempty_only, exclude_sundays)
            if factor and np.isfinite(factor):
                idx = act_pm["Month"] == forecast_last_ym
                if idx.any():
                    act_pm.loc[idx, "AfterVAT_excl_CC"] = act_pm.loc[idx, "AfterVAT_excl_CC"] * factor

    overlap = sorted(set(act_pm["Month"]).intersection(set(prev_pm["Month"])))
    if overlap:
        cmp_df = (pd.DataFrame({"Month": overlap})
                  .merge(act_pm, on="Month", how="left")
                  .merge(prev_pm.rename(columns={"AfterVAT_excl_CC":"Prev_AfterVAT_excl_CC"}), on="Month", how="left"))
        cmp_df = cmp_df.fillna(0.0)
        cmp_df["Delta"] = cmp_df["AfterVAT_excl_CC"] - cmp_df["Prev_AfterVAT_excl_CC"]
        cmp_df["% vs Prev"] = np.where(cmp_df["Prev_AfterVAT_excl_CC"] != 0,
                                       (cmp_df["Delta"]/cmp_df["Prev_AfterVAT_excl_CC"])*100.0, np.nan)
        df_cmp = cmp_df.sort_values("Month")
    else:
        df_cmp = pd.DataFrame(columns=["Month","AfterVAT_excl_CC","Prev_AfterVAT_excl_CC","Delta","% vs Prev"])

    total_actual_excl = float(pd.to_numeric(act_pm["AfterVAT_excl_CC"], errors="coerce").sum()) if not act_pm.empty else 0.0
    total_prev_excl   = float(pd.to_numeric(prev_pm["AfterVAT_excl_CC"], errors="coerce").sum()) if not prev_pm.empty else 0.0

    total_actual_incl = total_actual_excl + sum(_calc_after_vat(_cc_amount_for_month(cc_actual_df, ym), vat_rate, vat_mode)
                                                for ym in actual_months)
    total_prev_incl   = total_prev_excl + sum(_calc_after_vat(_cc_amount_for_month(cc_prev_df, ym), vat_rate, vat_mode)
                                              for ym in previous_months)

    df_tot = pd.DataFrame([
        {"Period":"Actual (Full Period)",   "Total_AfterVAT":round(total_actual_excl,2), "Total_AfterVAT_incl_CC":round(total_actual_incl,2)},
        {"Period":"Previous (Full Period)", "Total_AfterVAT":round(total_prev_excl,2),   "Total_AfterVAT_incl_CC":round(total_prev_incl,2)},
    ])
    return df_cmp, df_tot

# Guaranteed run_analysis (prefer core, else wrapper using core’s process/export)
if run_analysis_core is not None:
    run_analysis = run_analysis_core
else:
    def run_analysis(
        in_files,
        prev_file: str | None = None,
        call_center: float = 0.0,
        admin_forecast: float = 0.0,
        vat_rate: float = VAT_RATE_DEFAULT,
        vat_mode: str = VAT_MODE_DEFAULT,
        month: str | None = None,
        forecast_nonempty_only: bool = True,
        no_exclude_sundays: bool = False,
        out_name: str | None = None,
        prev_month_override: str | None = None,
    ) -> str:
        if process_dataframe_core is None or export_excel_core is None:
            raise ImportError("Core does not expose run_analysis nor process_dataframe/export_excel.")
        if not in_files:
            raise ValueError("No input files provided.")

        frames = []
        for p in in_files:
            df = read_excel_any(p)
            df["__source_file"] = os.path.basename(p)
            frames.append(df)
        df_all = pd.concat(frames, ignore_index=True)

        tables = process_dataframe_core(
            df_all,
            float(call_center), float(admin_forecast),
            float(vat_rate), str(vat_mode).lower(),
            SPECIAL_CORR_DEFAULT, CORRESPONDENT_MAP_DEFAULT,
            month, DATE_SOURCE_DEFAULT, TOP_N_NAMES_DEFAULT,
            bool(forecast_nonempty_only), not bool(no_exclude_sundays),
        )

        base  = out_name or f"monthly_revenue_VAT_{str(vat_mode).lower()}_{month or 'ALL'}"
        stamp = dt.datetime.now().strftime("%Y%m%d-%H%M")
        out_path = os.path.join(os.getcwd(), f"{base}__{stamp}.xlsx")
        export_excel_core(tables, out_path)
        return out_path

# ============================== UI ===============================
st.set_page_config(page_title="Artel Financial Suite — Comparisons", page_icon="📊", layout="wide")
st.title("📊 Artel Financial Suite — Range & Yearly Comparisons")

st.sidebar.header("⚙️ Settings")
vat_mode = st.sidebar.selectbox("VAT Mode (Revenue SAP)", ["extract", "add"], index=0)
vat_rate = st.sidebar.number_input("VAT rate (Revenue)", min_value=0.0, max_value=1.0, value=float(VAT_RATE_DEFAULT), step=0.01)
date_sources_all = ["Data of Document", "Data of transaction"]
date_sources = st.sidebar.multiselect("Date Source(s) for Revenue", date_sources_all, default=date_sources_all)
nonempty_only = st.sidebar.checkbox("Forecast: use only non-empty days", value=True)
exclude_sundays = st.sidebar.checkbox("Forecast: exclude Sundays", value=True)
st.sidebar.caption("SPECIAL_CORR is applied in normalization (G1 'ВЫЗОВ' exception).")

tab_rev, tab_exp, tab_cmp, tab_export = st.tabs(["📊 Revenue","💸 Expenditures (Yearly)","🧮 Comparison (Ranges)","📤 Export"])

for key in ["rev_df","exp_df","cc_actual_df","cc_prev_df","rev_cmp","rev_tot","exp_cmp","exp_tot"]:
    if key not in st.session_state: st.session_state[key] = None

# ---------------- Revenue ----------------
with tab_rev:
    st.subheader("Revenue Upload (SAP)")
    rev_files = st.file_uploader("Upload one or more SAP Revenue files (.xlsx/.xls/.csv)",
                                 type=["xlsx","xls","csv"], accept_multiple_files=True, key="rev_up")
    if rev_files:
        try:
            dfs = []
            for f in rev_files:
                cache_key = f"revcache::{f.name}::{len(f.getvalue())}"
                df = st.session_state.get(cache_key)
                if df is None:
                    df = _read_any_cached(f.getvalue(), f.name)
                    st.session_state[cache_key] = df
                dfs.append(df)
            d = pd.concat(dfs, ignore_index=True)

            v = validate_revenue(d)
            if not v.ok:
                st.error("⚠️ Invalid Revenue file structure.")
                st.write("Missing required columns:", v.missing)
                if v.suggestions: st.info(f"Suggestions: {v.suggestions}")
            else:
                st.success("✅ Revenue file(s) validated.")
                st.session_state.rev_df = normalize_revenue(d, corr_map=CORRESPONDENT_MAP_DEFAULT, special_corr=Special_CORR_DEFAULT if 'Special_CORR_DEFAULT' in globals() else SPECIAL_CORR_DEFAULT)
                st.dataframe(st.session_state.rev_df.head(20), use_container_width=True)
        except Exception as e:
            st.exception(e)

    st.markdown("---")
    st.subheader("Call Center & Admin (Month files)")
    c1, c2 = st.columns(2)
    with c1:
        cc_actual = st.file_uploader("Call Center - Actual Period (Month | Amount_USD)", type=["xlsx","csv"], key="cc_a")
        if cc_actual:
            try:
                if load_month_amount_file_core:
                    df = load_month_amount_file_core(io.BytesIO(cc_actual.getvalue()))
                else:
                    df = load_month_amount_file_local(cc_actual.getvalue())
                st.session_state.cc_actual_df = df
                st.caption("Parsed CC (Actual):"); st.dataframe(df, use_container_width=True)
            except Exception as e:
                st.exception(e)
    with c2:
        cc_prev = st.file_uploader("Call Center - Previous Period (Month | Amount_USD)", type=["xlsx","csv"], key="cc_p")
        if cc_prev:
            try:
                if load_month_amount_file_core:
                    df = load_month_amount_file_core(io.BytesIO(cc_prev.getvalue()))
                else:
                    df = load_month_amount_file_local(cc_prev.getvalue())
                st.session_state.cc_prev_df = df
                st.caption("Parsed CC (Previous):"); st.dataframe(df, use_container_width=True)
            except Exception as e:
                st.exception(e)

# --------------- Expenditures ----------------
with tab_exp:
    st.subheader("Expenditures Upload (RU headers)")
    exp_file = st.file_uploader("Upload Expenditures file (.xlsx/.xls/.csv) with RU headers",
                                type=["xlsx","xls","csv"], key="exp_up")
    if exp_file:
        try:
            cache_key = f"expcache::{exp_file.name}::{len(exp_file.getvalue())}"
            d = st.session_state.get(cache_key)
            if d is None:
                d = _read_any_cached(exp_file.getvalue(), exp_file.name)
                st.session_state[cache_key] = d

            if validate_expenditures is None or normalize_expenditures is None:
                st.error("Your core is missing validate_expenditures/normalize_expenditures.")
            else:
                v = validate_expenditures(d)
                if not v.ok:
                    st.error("⚠️ Invalid Expenditures file.")
                    st.write("Missing required columns (RU→EN):", v.missing)
                else:
                    st.success("✅ Expenditures file validated.")
                    st.session_state.exp_df = normalize_expenditures(d)
                    st.dataframe(st.session_state.exp_df.head(20), use_container_width=True)
        except Exception as e:
            st.exception(e)

    st.markdown("### Yearly Ranges")
    c1, c2 = st.columns(2)
    with c1:
        a_start = st.text_input("Actual Start (YYYY-MM)", "2025-01")
        a_end   = st.text_input("Actual End (YYYY-MM)",   "2025-12")
    with c2:
        p_start = st.text_input("Previous Start (YYYY-MM)", "2024-01")
        p_end   = st.text_input("Previous End (YYYY-MM)",   "2024-12")

    if st.button("Compare Expenditures (Yearly)"):
        if st.session_state.exp_df is None:
            st.warning("Upload Expenditures first.")
        elif compare_expenditures is None:
            st.error("Your core is missing compare_expenditures.")
        else:
            try:
                exp_cmp, exp_tot = compare_expenditures(
                    st.session_state.exp_df, a_start.strip(), a_end.strip(), p_start.strip(), p_end.strip()
                )
                if (exp_cmp is None or exp_cmp.empty) and (exp_tot is None or exp_tot.empty):
                    st.error("No rows matched the selected ranges. Check that 'Месяц' values are parseable (e.g., 'Январь 2025').")
                else:
                    st.subheader("By Category Comparison (Actual vs Previous)")
                    st.dataframe(exp_cmp, use_container_width=True) if (exp_cmp is not None and not exp_cmp.empty) else st.info("No category comparison to display.")
                    st.subheader("Totals")
                    st.dataframe(exp_tot, use_container_width=True) if (exp_tot is not None and not exp_tot.empty) else st.info("No totals to display.")
                    st.session_state.exp_cmp, st.session_state.exp_tot = exp_cmp, exp_tot
            except Exception as e:
                st.exception(e)

# --------------- Revenue range comparison ---------------
with tab_cmp:
    st.subheader("Revenue Range Comparison")
    c1, c2 = st.columns(2)
    with c1:
        ar_start = st.text_input("Actual Period Start (YYYY-MM)", "2025-01", key="ar_s")
        ar_end   = st.text_input("Actual Period End (YYYY-MM)",   "2025-10", key="ar_e")
    with c2:
        pr_start = st.text_input("Previous Period Start (YYYY-MM)", "2024-01", key="pr_s")
        pr_end   = st.text_input("Previous Period End (YYYY-MM)",   "2024-12", key="pr_e")

    forecast_last = st.checkbox("Forecast Actual End Month (active-days)", value=True)

    if st.button("Build Range Comparison (Revenue)"):
        if st.session_state.rev_df is None:
            st.warning("Upload Revenue first.")
        else:
            date_cols = [c for c in date_sources if c in st.session_state.rev_df.columns]
            if not date_cols:
                st.error("Selected date source(s) not found in data.\n\n"
                         f"Chosen: {date_sources}\nAvailable: {', '.join(map(str, st.session_state.rev_df.columns))}")
            else:
                try:
                    df_cmp, df_tot = compare_ranges_revenue(
                        st.session_state.rev_df,
                        ar_start.strip(), ar_end.strip(),
                        pr_start.strip(), pr_end.strip(),
                        date_cols,
                        float(vat_rate), str(vat_mode),
                        st.session_state.cc_actual_df, st.session_state.cc_prev_df,
                        (ar_end.strip() if forecast_last else None),
                        bool(nonempty_only), bool(exclude_sundays)
                    )
                    if (df_cmp is None or df_cmp.empty) and (df_tot is None or df_tot.empty):
                        st.error("Revenue comparison returned no rows. Check date ranges and that your files contain those months.")
                    else:
                        st.subheader("Overlap Comparison (After VAT, excl CC)")
                        st.dataframe(df_cmp, use_container_width=True) if (df_cmp is not None and not df_cmp.empty) else st.info("No overlap rows to display.")
                        st.subheader("Full-Period Totals")
                        st.dataframe(df_tot, use_container_width=True) if (df_tot is not None and not df_tot.empty) else st.info("No totals to display.")
                        st.session_state.rev_cmp, st.session_state.rev_tot = df_cmp, df_tot
                except Exception as e:
                    st.exception(e)

# ---------------------------- Export ----------------------------
with tab_export:
    st.subheader("Export Options")
    export_choice = st.radio("Select Export Type", [
        "Full Report (Revenue + Expenditures + Combined)",
        "Revenue Only", "Expenditures Only", "Combined Summary",
    ])
    report_title = st.text_input("Report Title", "Artel Financial Overview")
    subtitle     = st.text_input("Subtitle",    "Generated by Artel Financial Suite")

    sheets: Dict[str, pd.DataFrame] = {}
    if st.session_state.rev_cmp is not None: sheets["Revenue_Overlap"] = st.session_state.rev_cmp
    if st.session_state.rev_tot is not None: sheets["Revenue_Totals"] = st.session_state.rev_tot
    if st.session_state.exp_cmp is not None: sheets["EXP_ByCategory_Compare"] = st.session_state.exp_cmp
    if st.session_state.exp_tot is not None: sheets["EXP_Totals"] = st.session_state.exp_tot

    # Combined quick view (best-effort)
    if st.session_state.get("rev_tot") is not None and st.session_state.get("exp_tot") is not None:
        try:
            rev_tot = st.session_state.rev_tot; exp_tot = st.session_state.exp_tot
            comb = pd.DataFrame({
                "Metric": [
                    "Revenue Total (Actual)", "Revenue Total (Prev)",
                    "Expenditures Total (Actual)", "Expenditures Total (Prev)"
                ],
                "Value": [
                    float(rev_tot.loc[0, "Total_AfterVAT"]) if not rev_tot.empty else 0.0,
                    float(rev_tot.loc[1, "Total_AfterVAT"]) if not rev_tot.empty else 0.0,
                    float(exp_tot.loc[exp_tot["Period"]=="Actual (Full Period)",   "Total_Amount_USD"].values[0]) if not exp_tot.empty else 0.0,
                    float(exp_tot.loc[exp_tot["Period"]=="Previous (Full Period)", "Total_Amount_USD"].values[0]) if not exp_tot.empty else 0.0,
                ]
            })
            sheets["Combined_Summary"] = comb
        except Exception:
            pass

    def _filter(choice: str, all_sheets: Dict[str, pd.DataFrame]) -> Dict[str, pd.DataFrame]:
        if choice == "Revenue Only":       return {k:v for k,v in all_sheets.items() if k.startswith("Revenue")}
        if choice == "Expenditures Only":  return {k:v for k,v in all_sheets.items() if k.startswith("EXP_")}
        if choice == "Combined Summary":   return {k:v for k,v in all_sheets.items() if k.startswith("Combined")}
        return all_sheets

    chosen = _filter(export_choice, sheets)

    c1, c2, c3 = st.columns(3)
    with c1:
        if st.button("⬇️ Export Excel (.xlsx)"):
            if not chosen:
                st.warning("Nothing to export yet. Build a comparison first.")
            else:
                try:
                    xbytes = export_excel_bytes(chosen)
                    st.download_button("Download Excel", data=xbytes, file_name="financial_report.xlsx",
                                       mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
                except Exception as e:
                    st.exception(e)
    with c2:
        if st.button("⬇️ Export PowerPoint (.pptx)"):
            if not chosen:
                st.warning("Nothing to export yet.")
            else:
                try:
                    if build_pptx_core is None:
                        # Minimal built-in PPTX (2 tables)
                        from pptx import Presentation
                        from pptx.util import Inches
                        prs = Presentation()
                        slide = prs.slides.add_slide(prs.slide_layouts[0])
                        slide.shapes.title.text = report_title
                        slide.placeholders[1].text = subtitle
                        added = 0
                        for name, df in chosen.items():
                            if added >= 2: break
                            slide = prs.slides.add_slide(prs.slide_layouts[5])
                            slide.shapes.title.text = name
                            rows, cols = (min(12, len(df)) + 1, min(6, df.shape[1]))
                            x, y, cx, cy = Inches(0.5), Inches(1.2), Inches(9.0), Inches(4.5)
                            table = slide.shapes.add_table(rows, cols, x, y, cx, cy).table
                            for j, col in enumerate(df.columns[:cols]): table.cell(0,j).text = str(col)
                            for i in range(rows - 1):
                                for j in range(cols):
                                    try: val = df.iloc[i, j]
                                    except Exception: val = ""
                                    table.cell(i+1, j).text = "" if pd.isna(val) else str(val)
                            added += 1
                        bio = io.BytesIO(); prs.save(bio); bio.seek(0); pptx_bytes = bio.getvalue()
                    else:
                        # Use core builder
                        pptx_bytes = build_pptx_core(report_title, subtitle, charts={}, tables=chosen)
                    if pptx_bytes:
                        st.download_button("Download PPTX", data=pptx_bytes, file_name="financial_report.pptx",
                                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
                except Exception as e:
                    st.exception(e)
    with c3:
        if st.button("⬇️ Export PDF (.pdf)"):
            if not chosen:
                st.warning("Nothing to export yet.")
            else:
                try:
                    if build_pdf_core is None:
                        from reportlab.lib.pagesizes import A4
                        from reportlab.pdfgen import canvas
                        from reportlab.lib.units import cm
                        bio = io.BytesIO(); c = canvas.Canvas(bio, pagesize=A4); w, h = A4
                        c.setFont("Helvetica-Bold", 18); c.drawString(2*cm, h-3*cm, report_title)
                        c.setFont("Helvetica", 12); c.drawString(2*cm, h-4*cm, subtitle); c.showPage()
                        for name, df in list(chosen.items())[:2]:
                            c.setFont("Helvetica-Bold", 14); c.drawString(2*cm, h-2*cm, name)
                            c.setFont("Helvetica", 9); y = h-3*cm
                            cols = list(map(str, df.columns[:6])); c.drawString(2*cm, y, " | ".join(cols)); y -= 0.5*cm
                            for i in range(min(12, len(df))):
                                vals = [str(df.iloc[i, j]) for j in range(min(6, df.shape[1]))]
                                c.drawString(2*cm, y, " | ".join(vals)); y -= 0.5*cm
                                if y < 3*cm: c.showPage(); y = h-3*cm
                            c.showPage()
                        c.save(); bio.seek(0); pdf_bytes = bio.getvalue()
                    else:
                        pdf_bytes = build_pdf_core(report_title, subtitle, chosen)
                    if pdf_bytes:
                        st.download_button("Download PDF", data=pdf_bytes, file_name="financial_report.pdf",
                                           mime="application/pdf")
                except Exception as e:
                    st.exception(e)

st.markdown("---")
st.caption("Upload Revenue/Expenditures in the first tabs, run the comparisons, then export. All computations guarded by buttons.")
